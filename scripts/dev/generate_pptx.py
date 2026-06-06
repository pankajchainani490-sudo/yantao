import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[2]
SOURCE_MD = PROJECT_ROOT / "ppt" / "final_ppt_content.md"
OUTPUT_PPTX = PROJECT_ROOT / "ppt" / "malicious-traffic-defense.pptx"


def parse_sections(markdown: str):
    parts = re.split(r"^##\s+", markdown, flags=re.MULTILINE)
    slides = []

    for chunk in parts[1:]:
        lines = chunk.splitlines()
        title = lines[0].strip()
        body = "\n".join(lines[1:]).strip()
        slides.append((title, body))

    return slides


def extract_section(body: str, heading: str) -> str:
    pattern = rf"^###\s+{re.escape(heading)}\s*$"
    match = re.search(pattern, body, flags=re.MULTILINE)
    if not match:
        return ""

    start = match.end()
    next_heading = re.search(r"^###\s+.+$", body[start:], flags=re.MULTILINE)
    end = start + next_heading.start() if next_heading else len(body)
    return body[start:end].strip()


def clean_lines(block: str):
    lines = []
    for raw in block.splitlines():
        line = raw.strip()
        if not line:
            continue
        line = re.sub(r"^[-*]\s+", "", line)
        line = re.sub(r"^\d+\.\s+", "", line)
        lines.append(line)
    return lines


def add_title(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.8), Inches(0.8))
    frame = box.text_frame
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(17, 35, 68)


def add_cover_slide(prs: Presentation, title: str, subtitle: str, notes: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(241, 246, 255)
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.0), Inches(1.3))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = "Microsoft YaHei"
    run.font.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = RGBColor(16, 44, 84)

    subtitle_box = slide.shapes.add_textbox(Inches(1.2), Inches(3.0), Inches(10.6), Inches(0.9))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = subtitle
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(59, 92, 142)

    footer_box = slide.shapes.add_textbox(Inches(7.8), Inches(6.4), Inches(4.1), Inches(0.6))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "姓名 / 学号 / 专业 / 指导老师"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(90, 110, 140)

    slide.notes_slide.notes_text_frame.text = notes


def add_content_slide(prs: Presentation, title: str, content_lines, notes: str, layout_hint: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()

    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.95))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(224, 236, 252)
    header.line.fill.background()

    add_title(slide, title)

    body_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(8.2), Inches(5.8))
    frame = body_box.text_frame
    frame.word_wrap = True

    first = True
    for line in content_lines[:12]:
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(20 if len(content_lines) <= 6 else 18)
        paragraph.font.color.rgb = RGBColor(38, 55, 88)
        paragraph.space_after = Pt(8)
        if not line.endswith(":"):
            paragraph.bullet = True

    side_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.15), Inches(1.35), Inches(3.2), Inches(4.95))
    side_box.fill.solid()
    side_box.fill.fore_color.rgb = RGBColor(244, 248, 255)
    side_box.line.color.rgb = RGBColor(185, 205, 235)

    side_text = slide.shapes.add_textbox(Inches(9.35), Inches(1.55), Inches(2.8), Inches(4.55))
    side_frame = side_text.text_frame
    p1 = side_frame.paragraphs[0]
    p1.text = "版式建议"
    p1.font.name = "Microsoft YaHei"
    p1.font.bold = True
    p1.font.size = Pt(16)
    p1.font.color.rgb = RGBColor(31, 79, 141)

    p2 = side_frame.add_paragraph()
    p2.text = layout_hint or "按正文要点排版即可"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(73, 91, 122)
    p2.space_before = Pt(8)

    note_tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.15), Inches(6.05), Inches(3.2), Inches(0.5))
    note_tag.fill.solid()
    note_tag.fill.fore_color.rgb = RGBColor(231, 241, 255)
    note_tag.line.fill.background()
    note_tf = note_tag.text_frame
    note_tf.paragraphs[0].text = "备注见讲稿"
    note_tf.paragraphs[0].font.name = "Microsoft YaHei"
    note_tf.paragraphs[0].font.size = Pt(12)
    note_tf.paragraphs[0].font.color.rgb = RGBColor(62, 94, 145)
    note_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    slide.notes_slide.notes_text_frame.text = notes


def build_pptx():
    markdown = SOURCE_MD.read_text(encoding="utf-8")
    slides = parse_sections(markdown)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for index, (title, body) in enumerate(slides):
        page_title = re.sub(r"^第\s*\d+\s*页\s*", "", title).strip()
        notes = extract_section(body, "演讲备注")
        layout_hint = extract_section(body, "版式建议")

        if index == 0:
            cover_title = extract_section(body, "标题").splitlines()[0].strip()
            subtitle_lines = clean_lines(extract_section(body, "副标题"))
            subtitle = subtitle_lines[0] if subtitle_lines else ""
            add_cover_slide(prs, cover_title, subtitle, notes)
            continue

        main_block = extract_section(body, "页面正文")
        lines = clean_lines(main_block)
        add_content_slide(prs, page_title, lines, notes, "\n".join(clean_lines(layout_hint)))

    prs.save(OUTPUT_PPTX)
    print(f"Generated PPTX: {OUTPUT_PPTX}")


if __name__ == "__main__":
    build_pptx()
